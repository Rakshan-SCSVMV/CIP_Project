import os
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    print("Library 'python-pptx' not found. Please run: pip install python-pptx")
    exit()

def create_presentation():
    prs = Presentation()
    
    # helper to set background color
    def set_bg(slide, r, g, b):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(r, g, b)

    # 1. Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide, 0, 47, 108) # Navy
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Digital Land Litigation Management System"
    subtitle.text = "DLLMS v2.0 | Next-Generation Government Portal\nPresented by: [Your Name]"
    
    # Style title
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(232, 160, 0) # Accent Gold

    # 2. Problem Statement
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "The Problem"
    content = slide.placeholders[1].text_frame
    content.text = "Existing challenges in Land Dispute Resolution:"
    points = [
        "Manual, paper-based workflows causing massive backlogs.",
        "Zero transparency for citizens during case lifecycles.",
        "Difficulties in cross-departmental record synchronization.",
        "Physical distance barriers for rural residents."
    ]
    for p in points:
        tf = content.add_paragraph()
        tf.text = p
        tf.level = 1

    # 3. Our Solution
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Our Solution: DLLMS v2.0"
    content = slide.placeholders[1].text_frame
    content.text = "A Unified, Citizen-Centric Digital Ecosystem:"
    points = [
        "End-to-End e-Filing & Tracking interface.",
        "Automated Real-time Analytics for decision making.",
        "Multilingual Portal support (English & Tamil).",
        "Zero-configuration Portable Backend architecture."
    ]
    for p in points:
        tf = content.add_paragraph()
        tf.text = p
        tf.level = 1

    # 4. System Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "System Architecture"
    content = slide.placeholders[1].text_frame
    content.text = "Technical Stack Overview:"
    points = [
        "Frontend: Modern HTML5, CSS3, ES6 JavaScript SPA.",
        "Backend: Professional Python API Server (Threaded).",
        "Data Store: NoSQL JSON Database (Location Agnostic).",
        "Utilities: Advanced Seeder for demo-ready high-quality data."
    ]
    for p in points:
        tf = content.add_paragraph()
        tf.text = p
        tf.level = 1

    # 5. Core Features: Citizen
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Citizen Empowerment"
    content = slide.placeholders[1].text_frame
    points = [
        "Smart Case Filing: Guided 5-Step intuitive process.",
        "Case Passport: 1-click status view with full legal history.",
        "Glossary & FAQs: Detailed legal process support.",
        "Integrated AI Helpdesk: 24x7 automated query resolution."
    ]
    for p in points:
        tf = content.add_paragraph()
        tf.text = p
        tf.level = 0

    # 6. Core Features: Admin
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Administrative Intelligence"
    content = slide.placeholders[1].text_frame
    points = [
        "Decision Dashboard: Real-time trends & dispute heatmaps.",
        "Officer Portal: Manage hearings and status updates.",
        "Advanced Search: Multi-filter engine for district-level tracking.",
        "Data Export: Standardized CSV reporting for auditing."
    ]
    for p in points:
        tf = content.add_paragraph()
        tf.text = p
        tf.level = 0

    # 7. UX & Aesthetics
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Design & User Experience"
    content = slide.placeholders[1].text_frame
    points = [
        "Elite Aesthetics: Navy & Accent Gold government theme.",
        "Rich Typography: Noto Sans & JetBrains Mono integration.",
        "Micro-animations: Responsive hover states and transitions.",
        "Accessibility: High-contrast mode and font scaling support."
    ]
    for p in points:
        tf = content.add_paragraph()
        tf.text = p
        tf.level = 0

    # 8. Live Demo Slides
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "System Real-Time Analytics"
    content = slide.placeholders[1].text_frame
    content.text = "[Insert Analytics Dashboard Screenshot Here]"
    
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Advanced Search & Filing"
    content = slide.placeholders[1].text_frame
    content.text = "[Insert Case Tracker Screenshot Here]"

    # 9. Future Scope
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Future Road Map"
    content = slide.placeholders[1].text_frame
    points = [
        "Cloud Integration: Scalable hosting on National NIC infrastructure.",
        "AI Engine: Automatic document sentiment & risk analysis.",
        "Mobile App: Dedicated citizen application for on-site filing.",
        "Interoperability: API connection with e-Courts and DILRMP."
    ]
    for p in points:
        tf = content.add_paragraph()
        tf.text = p
        tf.level = 0

    # Slide 11: Future Roadmap / v3.0 Features
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "New in v3.0: National Scale"
    tf = slide.placeholders[1].text_frame
    tf.text = "Expanded Features for Government-Grade Utility"
    tf.add_paragraph().text = "• 26,000+ Supreme Court Judgments Integration"
    tf.add_paragraph().text = "• National Geographic Cascading Search (State/District/Village)"
    tf.add_paragraph().text = "• High-Performance Metadata Indexing"
    tf.add_paragraph().text = "• Ready for Large-Scale Deployment"

    # Slide 12: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Thank You"
    slide.placeholders[1].text = "Questions & Discussion\nContact: support@dllms.gov.in"

    file_path = "DLLMS_Presentation_v3.pptx"
    prs.save(file_path)
    print(f"Presentation created successfully: {file_path}")

if __name__ == "__main__":
    create_presentation()
