
from pptx import Presentation

def inspect_pptx(path):
    try:
        prs = Presentation(path)
        print(f"Inspecting: {path}")
        for i, slide in enumerate(prs.slides):
            title = "No Title"
            if slide.shapes.title:
                title = slide.shapes.title.text
            print(f"Slide {i+1}: {title}")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    print(f"  - {shape.text[:100]}")
    except Exception as e:
        print(f"Error: {e}")

inspect_pptx(r"c:\master project\LLS\LLS\ppt instruction\CIP   PPT 25-26.pptx")
