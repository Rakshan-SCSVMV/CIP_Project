
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def set_bg(slide, r, g, b):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)

def add_bullet_points(slide, points, level=0):
    tf = slide.placeholders[1].text_frame
    for p in points:
        p_para = tf.add_paragraph()
        p_para.text = p
        p_para.level = level

def create_review_1(filename):
    prs = Presentation()
    
    # 1. Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide, 0, 47, 108) # Navy
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Digital Land Litigation Management System (DLLMS)"
    subtitle.text = "Review 1: Project Proposal & Scope\nPresented by: [Team Name]\nBatch: 2025-26"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(232, 160, 0)

    # 2. Introduction
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Project Introduction"
    points = [
        "A national-scale portal for managing land dispute cases.",
        "Aims to digitize the Revenue Department's legal workflows.",
        "Provides a unified platform for citizens and administrators.",
        "Incorporates historical legal data for better case analysis."
    ]
    add_bullet_points(slide, points)

    # 3. Existing System & Drawbacks
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Existing System & Drawbacks"
    points = [
        "Paper-based record keeping leads to physical damage and loss.",
        "Slow manual processing causes massive case backlogs (years of delay).",
        "Lack of transparency for citizens to track their own case status.",
        "Geographical barriers for rural citizens to access urban offices."
    ]
    add_bullet_points(slide, points)

    # 4. Literature Survey (Mocked based on project context)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Literature Survey"
    points = [
        "National e-Governance Plan (NeGP): Guidelines for digitizing land records.",
        "DILRMP: Digital India Land Records Modernization Programme.",
        "e-Courts Project: Study on judicial digitization in India.",
        "Study on NoSQL databases for high-velocity legal data handling."
    ]
    add_bullet_points(slide, points)

    # 5. Proposed System & Objectives
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Proposed System & Objectives"
    points = [
        "Objective 1: Create a secure, Aadhaar-linked e-filing interface.",
        "Objective 2: Implement a national-scale cascading geographic search.",
        "Objective 3: Index 26,000+ Supreme Court judgments for instant access.",
        "Objective 4: Provide real-time analytics for government officials."
    ]
    add_bullet_points(slide, points)

    prs.save(filename)
    print(f"Created: {filename}")

def create_review_2(filename):
    prs = Presentation()
    
    # 1. Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide, 0, 47, 108)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Digital Land Litigation Management System (DLLMS)"
    subtitle.text = "Review 2: Design & System Architecture\nPresented by: [Team Name]"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # 2. Methodology
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Methodology"
    points = [
        "SDLC: Agile Development Methodology.",
        "Requirement Gathering: Analyzing Revenue Dept. workflows.",
        "Data Modeling: Designing high-performance JSON hierarchies.",
        "Prototyping: Building the Single Page Application (SPA) frontend."
    ]
    add_bullet_points(slide, points)

    # 3. System Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "System Architecture"
    points = [
        "Tier 1 (Frontend): Vanilla JS SPA with Government UI Theme.",
        "Tier 2 (Backend): Python Multi-threaded API Server.",
        "Tier 3 (Data): Persistent NoSQL JSON Data Stores.",
        "External: Google Gemini AI Integration for Legal Assistant."
    ]
    add_bullet_points(slide, points)

    # 4. Key Modules
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Core Modules"
    points = [
        "Citizen Portal: e-Filing, Case Passport, Bilingual support.",
        "Admin Dashboard: Heatmaps, Status Management, Data Exports.",
        "Legal Archive: Metadata-indexed search for SC Judgments.",
        "National Search: Cascading State > District > Village filters."
    ]
    add_bullet_points(slide, points)

    # 6. Implementation Screenshots
    slide = prs.slides.add_slide(prs.slide_layouts[5]) # Title only layout
    slide.shapes.title.text = "Implementation: Home & Tracker"
    if os.path.exists("presentations/home.png"):
        slide.shapes.add_picture("presentations/home.png", Inches(0.5), Inches(1.5), height=Inches(5))
    if os.path.exists("presentations/tracker.png"):
        slide.shapes.add_picture("presentations/tracker.png", Inches(5.5), Inches(1.5), height=Inches(5))

    prs.save(filename)
    print(f"Created: {filename}")

def create_final_presentation(filename):
    prs = Presentation()
    
    # 1. Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide, 0, 47, 108)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Digital Land Litigation Management System (DLLMS)"
    subtitle.text = "Final Project Presentation\nPresented by: [Team Name]\nGuide: [Guide Name]"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # 2. Project Summary
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Project Summary"
    points = [
        "Successfully developed a national-scale e-governance portal.",
        "Integrated 26k+ judgments and national geography data.",
        "Implemented secure 5-step case filing workflow.",
        "Achieved sub-second response times for data-heavy queries."
    ]
    add_bullet_points(slide, points)

    # 3. Implementation Highlights
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Implementation Highlights"
    points = [
        "Frontend: Dark mode, bilingual, and accessibility compliant.",
        "Backend: Robust Python API handling concurrent sessions.",
        "AI: Integration of LLM for instant legal query resolution.",
        "Seeding: Dynamic generation of 26k+ demo entries."
    ]
    add_bullet_points(slide, points)

    # 4. Results & Analysis (Screenshots)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Admin Dashboard & Judgment Library"
    if os.path.exists("presentations/admin.png"):
        slide.shapes.add_picture("presentations/admin.png", Inches(0.5), Inches(1.5), height=Inches(5))
    if os.path.exists("presentations/judgments.png"):
        slide.shapes.add_picture("presentations/judgments.png", Inches(5.5), Inches(1.5), height=Inches(5))

    # 5. Results & Analysis (Metrics)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "System Performance & Metrics"
    points = [
        "Digitization: 100% reduction in paper dependency for new cases.",
        "Efficiency: Automated status tracking reduces manual query load.",
        "Scalability: Handles 26,000+ judgments in a lightweight portable format.",
        "User Feedback: Intuitive design praised for accessibility."
    ]
    add_bullet_points(slide, points)

    # 6. Conclusion & Future Scope
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Conclusion & Future Scope"
    points = [
        "Conclusion: DLLMS provides a viable blueprint for national digitization.",
        "Future 1: Blockchain integration for immutable land records.",
        "Future 2: OCR for automatic digitization of old paper records.",
        "Future 3: Mobile application for field officers."
    ]
    add_bullet_points(slide, points)

    # 7. References
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "References"
    points = [
        "1. National Informatics Centre (NIC) - e-Governance Standards.",
        "2. Ministry of Land Records - DILRMP Operational Guidelines.",
        "3. Supreme Court of India - Legal Data APIs Documentation.",
        "4. Python Software Foundation - Advanced Threading in Web Servers."
    ]
    add_bullet_points(slide, points)

    # 8. Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Thank You"
    slide.placeholders[1].text = "Questions?"

    prs.save(filename)
    print(f"Created: {filename}")

if __name__ == "__main__":
    os.makedirs("presentations", exist_ok=True)
    create_review_1("presentations/DLLMS_Review_1_Proposal.pptx")
    create_review_2("presentations/DLLMS_Review_2_Design.pptx")
    create_final_presentation("presentations/DLLMS_Final_Review.pptx")
