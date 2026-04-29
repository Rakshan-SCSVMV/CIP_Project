from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def set_bg(slide, r, g, b):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)

def add_title_slide(prs, title_text, subtitle_text, r=0, g=47, b=108):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide, r, g, b)
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    subtitle = slide.placeholders[1]
    subtitle.text = subtitle_text
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

def add_content_slide(prs, title_text, bullet_points):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.text = bullet_points[0]
    for bp in bullet_points[1:]:
        p = tf.add_paragraph()
        p.text = bp
        p.level = 0

def create_review_0():
    prs = Presentation()
    add_title_slide(prs, "Review 0: Project Initiation", "Digital Land Litigation Management System (DLLMS)\nDate: 18.02.2026")
    add_content_slide(prs, "Domain Identification", ["• Domain: e-Governance & Public Administration", "• Sub-Domain: Land Revenue & Judicial Records", "• Focus: Digital transformation of paper-based litigation"])
    add_content_slide(prs, "Project Abstract", ["• A national-scale portal for managing land disputes.", "• Features: Case filing, tracking, and legal archive access.", "• Goal: Transparency, speed, and accessibility in revenue courts."])
    prs.save("Review_0.pptx")

def create_review_1():
    prs = Presentation()
    add_title_slide(prs, "Review 1: System Requirements & Architecture", "Digital Land Litigation Management System\nDate: 06.03.2026")
    add_content_slide(prs, "Problem Statement", ["• Manual tracking leads to years of delays.", "• Lack of a unified national database for land disputes.", "• Difficulty in accessing legal precedents (Supreme Court Judgments)."])
    add_content_slide(prs, "Proposed System", ["• Web-based SPA (Single Page Application) for real-time tracking.", "• Integrated National Geographic Search (Cascading).", "• AI-powered legal assistant for citizen guidance."])
    add_content_slide(prs, "System Architecture", ["• Frontend: HTML5, Vanilla CSS3, JavaScript (ES8+).", "• Backend: Threaded Python API Server.", "• Data: JSON-based NoSQL for high-speed local deployment."])
    prs.save("Review_1.pptx")

def create_review_2():
    prs = Presentation()
    add_title_slide(prs, "Review 2: Implementation & 50% Coding", "Digital Land Litigation Management System\nDate: 02.04.2026")
    add_content_slide(prs, "Implementation Status (50%)", ["• UI/UX Design completed with Government-style aesthetics.", "• Core Backend APIs for Case Management implemented.", "• National Geography Database (State/District/Village) integrated."])
    add_content_slide(prs, "UML & Techniques", ["• Techniques: High-speed metadata indexing.", "• Tools: Ripgrep for judgment archive searching.", "• Algorithm: Cascading filter logic for geographic precision."])
    prs.save("Review_2.pptx")

def create_review_3_4():
    prs = Presentation()
    add_title_slide(prs, "Review 3 & 4: Full Demo & Report", "Digital Land Litigation Management System\nDate: 23.04.2026 - 04.05.2026")
    add_content_slide(prs, "Complete Module Demo", ["• National Search Portal: Find cases by any location in India.", "• Judgment Library: Access 26k+ Supreme Court records.", "• AI Assistant: Real-time chat powered by Gemini simulation."])
    add_content_slide(prs, "System Highlights", ["• Click sounds and notification center for high engagement.", "• Full Dossier view for legal records (CNR, Advocates, etc.).", "• Print-ready official reports."])
    prs.save("Review_3_4.pptx")

def create_review_final():
    prs = Presentation()
    # 100% Professional Presentation
    add_title_slide(prs, "Final Presentation: DLLMS v3.0", "National Land Litigation Management System\nProject University Exam - May 2026", 0, 47, 108)
    
    add_content_slide(prs, "Executive Summary", ["• DLLMS is a comprehensive solution for India's Land Revenue Department.", "• Digitizes 26,000+ legal records and thousands of active cases.", "• Provides a transparent, AI-driven interface for citizens and officers."])
    
    add_content_slide(prs, "Key Modules & Innovation", ["• AI Chatbot: Simulated Gemini 1.5 Pro for legal guidance.", "• Geo-Portal: Cascading search for 28 States & 8 UTs.", "• Legal Dossier: Professional CNR-based tracking system."])
    
    add_content_slide(prs, "Technical Excellence", ["• Performance: JSON data handling for millisecond response times.", "• Accessibility: Bilingual support and High Contrast modes.", "• Scalability: Portable architecture ready for national rollout."])
    
    add_content_slide(prs, "Project Outcome", ["• 100% Implementation completed.", "• Professional UI with government-grade aesthetics.", "• Full support for physical and digital report submission."])
    
    prs.save("Review_Final.pptx")

if __name__ == "__main__":
    create_review_0()
    create_review_1()
    create_review_2()
    create_review_3_4()
    create_review_final()
    print("All Review PPTs created successfully.")
